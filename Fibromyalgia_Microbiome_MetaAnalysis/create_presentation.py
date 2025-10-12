#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation from scientific committee presentation markdown
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def parse_markdown_slides(markdown_content):
    """Parse markdown content and extract slide structure"""
    slides = []
    current_slide = None
    in_content = False
    table_rows = []
    in_table = False

    lines = markdown_content.split('\n')

    for line in lines:
        line = line.strip()

        # Skip empty lines
        if not line:
            continue

        # Slide title (starts with # or ---)
        if line.startswith('## ') and '##' in line:
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': line.replace('## ', '').split(':')[0].strip(), 'content': []}
            in_content = True
            table_rows = []
            in_table = False
        elif line.startswith('# ') and '#' in line and not current_slide:
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': line.replace('# ', '').split(':')[0].strip(), 'content': []}
            in_content = True
            table_rows = []
            in_table = False
        elif line == '---' or line.startswith('***') or line.startswith('___'):
            if current_slide:
                slides.append(current_slide)
            current_slide = None
            in_content = False
            table_rows = []
            in_table = False
        elif current_slide and (line.startswith('- ') or line.startswith('• ') or re.match(r'^\d+\.', line)):
            # Bullet point
            current_slide['content'].append({'type': 'bullet', 'text': line})
        elif current_slide and '|' in line and line.startswith('|') and not line.startswith('|---'):
            # Table row
            table_rows.append(line)
            if not in_table:
                in_table = True
        elif current_slide and line.startswith('|---|---') and in_table:
            # Table separator - process complete table
            headers = [col.strip() for col in table_rows[0].split('|')[1:-1]]
            data_rows = []
            for row in table_rows[1:]:
                if '|' in row:
                    cols = [col.strip() for col in row.split('|')[1:-1]]
                    data_rows.append(cols)
            current_slide['content'].append({'type': 'table', 'headers': headers, 'rows': data_rows})
            table_rows = []
            in_table = False
        elif current_slide and line.startswith('`'):
            # Code block
            current_slide['content'].append({'type': 'code', 'text': line.strip('`')})
        elif current_slide and not in_table:
            # Regular text
            current_slide['content'].append({'type': 'text', 'text': line})

    # Add the last slide
    if current_slide:
        slides.append(current_slide)

    return slides

def create_slide_title_only(prs, title):
    """Create a title-only slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    return slide

def create_slide_title_and_content(prs, title, content_items):
    """Create a title and content slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    content_placeholder = slide.placeholders[1]

    # Clear existing content
    content_placeholder.text = ''

    for item in content_items:
        if item['type'] == 'bullet':
            p = content_placeholder.text_frame.add_paragraph()
            p.text = item['text'].lstrip('- •').strip()
            p.level = 0
        elif item['type'] == 'text':
            p = content_placeholder.text_frame.add_paragraph()
            p.text = item['text']
            p.level = 0
        elif item['type'] == 'code':
            p = content_placeholder.text_frame.add_paragraph()
            p.text = f"```{item['text']}```"
            p.level = 0

    return slide

def format_slide_text(slide):
    """Apply consistent formatting to slide text"""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
                    run.font.name = 'Arial'
                paragraph.alignment = PP_ALIGN.LEFT

def create_powerpoint_from_markdown(markdown_file_path, output_path):
    """Main function to create PPT from markdown"""
    # Read markdown content
    with open(markdown_file_path, 'r', encoding='utf-8') as f:
        markdown_content = f.read()

    # Parse slides
    slides_data = parse_markdown_slides(markdown_content)

    # Create presentation
    prs = Presentation()

    # Title slide
    title_slide = create_slide_title_only(prs, "Fibromyalgia-Microbiome Diversity Meta-Analysis")
    format_slide_text(title_slide)

    # Add content slides
    for slide_data in slides_data[1:]:  # Skip title slide
        slide = create_slide_title_and_content(
            prs,
            slide_data['title'],
            slide_data.get('content', [])
        )
        format_slide_text(slide)

    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

def create_multiple_presentations():
    """Create presentations from multiple markdown sources"""

    presentations = [
        ('scientific_committee_presentation.md', 'scientific_committee_presentation.pptx'),
        ('PUBLICATION_READY_FINAL_MANUSCRIPT.md', 'final_manuscript_presentation.pptx'),
        ('supplementary_materials_complete.md', 'supplementary_materials_presentation.pptx')
    ]

    for markdown_path, output_path in presentations:
        if os.path.exists(markdown_path):
            print(f"Creating presentation from: {markdown_path}")
            create_powerpoint_from_markdown(markdown_path, output_path)
            print(f"✓ Created: {output_path}")
        else:
            print(f"✗ Markdown file not found: {markdown_path}")

if __name__ == "__main__":
    print("Creating multiple PowerPoint presentations from markdown sources...")
    create_multiple_presentations()
    print("\nAll presentations created successfully!")
